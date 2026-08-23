"""模块 2.4 课件文本提取测试。"""

from pathlib import Path

from app.services.courseware import extract_courseware


def test_extract_markdown(tmp_path):
    f = tmp_path / "c.md"
    f.write_text("# 第一章\n内容A\n## 1.1 小节\n内容B\n", encoding="utf-8")
    text, has_chapters, warning = extract_courseware(f, "md")
    assert "内容A" in text
    assert has_chapters is True
    assert warning is None


def test_extract_markdown_no_chapters(tmp_path):
    f = tmp_path / "c.md"
    f.write_text("纯文本没有标题\n", encoding="utf-8")
    _, has_chapters, _ = extract_courseware(f, "md")
    assert has_chapters is False


def test_extract_pdf(tmp_path):
    import pymupdf

    f = tmp_path / "c.pdf"
    doc = pymupdf.open()
    page = doc.new_page()
    page.insert_text((72, 72), "Chapter 1 Test Content")
    doc.save(str(f))
    doc.close()

    text, has_chapters, warning = extract_courseware(f, "pdf")
    assert "Test Content" in text
    assert warning is None


def test_extract_pptx(tmp_path):
    from pptx import Presentation

    f = tmp_path / "c.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "第一章 标题"
    prs.save(str(f))

    text, has_chapters, warning = extract_courseware(f, "pptx")
    assert "标题" in text
    assert has_chapters is True
    assert warning is None


def test_extract_pptx_empty_warning(tmp_path):
    from pptx import Presentation

    f = tmp_path / "empty.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白版式
    # 移除所有 shape（模拟纯图片型 PPT，无任何文本）
    for shape in list(slide.shapes):
        shape._element.getparent().remove(shape._element)
    prs.save(str(f))

    text, _, warning = extract_courseware(f, "pptx")
    assert warning is not None
    assert "图片型" in warning


def test_extract_unsupported_format(tmp_path):
    f = tmp_path / "c.txt"
    f.write_text("x")
    try:
        extract_courseware(f, "txt")
        assert False, "应抛出异常"
    except ValueError:
        pass
