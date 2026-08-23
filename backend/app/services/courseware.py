"""课件文本提取（md / pdf / pptx）。

返回 (文本, 是否有章节结构, warning)。warning 为 None 表示正常。
"""

from pathlib import Path


def _extract_markdown(path: Path) -> tuple[str, bool, str | None]:
    text = path.read_text(encoding="utf-8")
    # 检测章节结构：标题行（# 开头）
    has_chapters = any(line.strip().startswith("#") for line in text.splitlines())
    return text, has_chapters, None


def _extract_pdf(path: Path) -> tuple[str, bool, str | None]:
    import pymupdf  # 延迟导入，避免未安装时影响其他功能

    doc = pymupdf.open(str(path))
    pages = []
    for page in doc:
        pages.append(page.get_text())
    doc.close()
    text = "\n\n".join(pages)
    # PDF 章节检测：常见章节标题模式（第X章/第X节/数字标题）
    import re

    has_chapters = bool(re.search(r"(第[一二三四五六七八九十\d]+[章节]|^\d+[\.、]\s*\S)", text, re.MULTILINE))
    return text, has_chapters, None


def _extract_pptx(path: Path) -> tuple[str, bool, str | None]:
    from pptx import Presentation  # 延迟导入

    prs = Presentation(str(path))
    parts = []
    has_chapters = False
    has_real_text = False
    for idx, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs).strip()
                    if line:
                        texts.append(line)
        # 第一段视为该页标题
        page_title = texts[0] if texts else ""
        if page_title:
            has_chapters = True
            parts.append(f"【第{idx}页 {page_title}】")
        else:
            parts.append(f"【第{idx}页】")
        body = texts[1:] if page_title else texts
        parts.extend(body)
        if texts:
            has_real_text = True

    text = "\n".join(parts)
    warning = None
    if not has_real_text:
        warning = "该 PPT 可能为纯图片型，未能提取到文本，建议导出为 PDF 后上传"
    return text, has_chapters, warning


def extract_courseware(path: Path, fmt: str) -> tuple[str, bool, str | None]:
    """按格式提取课件文本。fmt ∈ {md, pdf, pptx}。"""
    fmt = fmt.lower()
    if fmt == "md":
        return _extract_markdown(path)
    if fmt == "pdf":
        return _extract_pdf(path)
    if fmt == "pptx":
        return _extract_pptx(path)
    raise ValueError(f"不支持的课件格式: {fmt}")
